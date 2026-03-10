from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


GRADE_COLORS = {
    "A": RGBColor(35, 133, 63),
    "B": RGBColor(201, 170, 35),
    "C": RGBColor(217, 119, 6),
    "D": RGBColor(191, 43, 40),
}


@dataclass(slots=True)
class LeadCard:
    company_name: str | None
    sender_name: str | None
    sender_title: str | None
    sender_email: str
    grade: str
    score: int
    why_hot: str
    company_context: str
    suggested_approach: str
    suggested_opener: str
    key_signals: list[str]
    use_case: str | None = None
    estimated_acv: int | None = None
    meddpicc_summary: str | None = None
    competitive_landscape: str | None = None
    top_discovery_question: str | None = None


@dataclass(slots=True)
class DailyBriefing:
    briefing_date: str
    total_new: int
    total_updated: int
    grade_distribution: dict[str, int]
    top_intents: dict[str, int]
    executive_summary: str
    pipeline_value: str
    quick_wins: list[str]
    needs_attention: list[str]
    top_leads: list[LeadCard]


def format_markdown_briefing(briefing: DailyBriefing) -> str:
    grade_parts = ", ".join(f"{k}: {v}" for k, v in briefing.grade_distribution.items())
    intent_parts = ", ".join(f"{k}: {v}" for k, v in briefing.top_intents.items())

    lines: list[str] = [
        f"# Daily Lead Briefing — {briefing.briefing_date}",
        "",
        "## Executive Summary",
        briefing.executive_summary,
        "",
        "## Pipeline Snapshot",
        f"- **Pipeline Value (A+B leads):** {briefing.pipeline_value}",
        f"- Total new: {briefing.total_new}",
        f"- Total updated: {briefing.total_updated}",
        f"- Grade distribution: {grade_parts}",
        f"- Intent categories: {intent_parts}",
        "",
        "## Top Priority Leads",
        "",
    ]

    for idx, lead in enumerate(briefing.top_leads, start=1):
        acv_str = f"${lead.estimated_acv:,}" if lead.estimated_acv else "TBD"
        use_case_str = (lead.use_case or "unknown").replace("_", " ").title()
        company = lead.company_name or "Unknown Company"

        li_query = f"{lead.sender_name or ''} {lead.company_name or ''}".strip().replace(" ", "+")
        li_url = f"https://www.linkedin.com/search/results/all/?keywords={li_query}"

        lines.extend([
            f"### {idx}. {company} — {lead.grade} / {lead.score} — ACV: {acv_str}",
            f"**Contact:** {lead.sender_name or 'Unknown'}"
            + (f", {lead.sender_title}" if lead.sender_title else "")
            + f" ({lead.sender_email})"
            + f" — [LinkedIn]({li_url})",
            f"**Use Case:** {use_case_str}",
            "",
            f"**Why hot:** {lead.why_hot}",
            "",
            f"**Company context:** {lead.company_context}",
            "",
        ])

        if lead.competitive_landscape:
            lines.append(f"**Competitive landscape:** {lead.competitive_landscape}")
            lines.append("")

        if lead.meddpicc_summary:
            lines.append(f"**MEDDPICC status:** {lead.meddpicc_summary}")
            lines.append("")

        if lead.top_discovery_question:
            lines.append(f"**Top discovery question:** {lead.top_discovery_question}")
            lines.append("")

        lines.extend([
            f"**Suggested approach:** {lead.suggested_approach}",
            "",
            f"**Suggested opener:** {lead.suggested_opener}",
            "",
            "**Key signals:**",
        ])
        for signal in lead.key_signals:
            lines.append(f"- {signal}")
        lines.append("")

    lines.extend(["## Quick Wins", ""])
    for item in briefing.quick_wins:
        lines.append(f"- {item}")

    lines.extend(["", "## Needs Attention", ""])
    for item in briefing.needs_attention:
        lines.append(f"- {item}")

    lines.append("")
    return "\n".join(lines)


def write_markdown_report(briefing: DailyBriefing, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(format_markdown_briefing(briefing), encoding="utf-8")
    return output_path


def build_pptx_deck(briefing: DailyBriefing, output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, briefing)
    _add_pipeline_slide(prs, briefing)
    for lead in briefing.top_leads[:5]:
        _add_lead_slide(prs, lead)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _add_title_slide(prs: Presentation, briefing: DailyBriefing) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(12), Inches(0.8))
    p = title.text_frame.paragraphs[0]
    if os.environ.get("GITLAB_MODE", "false").lower() == "true":
        p.text = "GitLab Sales — Daily Lead Briefing"
    else:
        p.text = "Anthropic Mid-Market — Daily Lead Briefing"
    p.font.size = Pt(28)
    p.font.bold = True

    date_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(0.5))
    p = date_box.text_frame.paragraphs[0]
    p.text = f"{briefing.briefing_date}  |  Pipeline: {briefing.pipeline_value}  |  New: {briefing.total_new}  |  Updated: {briefing.total_updated}"
    p.font.size = Pt(14)

    summary = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(11.8), Inches(3.0))
    tf = summary.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = briefing.executive_summary
    p.font.size = Pt(16)

    quick = slide.shapes.add_textbox(Inches(0.7), Inches(5.3), Inches(5.8), Inches(1.8))
    tf = quick.text_frame
    p = tf.paragraphs[0]
    p.text = "Quick Wins"
    p.font.size = Pt(15)
    p.font.bold = True
    for item in briefing.quick_wins[:3]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)

    attention = slide.shapes.add_textbox(Inches(6.8), Inches(5.3), Inches(5.8), Inches(1.8))
    tf = attention.text_frame
    p = tf.paragraphs[0]
    p.text = "Needs Attention"
    p.font.size = Pt(15)
    p.font.bold = True
    for item in briefing.needs_attention[:3]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)


def _add_pipeline_slide(prs: Presentation, briefing: DailyBriefing) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
    p = title.text_frame.paragraphs[0]
    p.text = "Pipeline Overview"
    p.font.size = Pt(24)
    p.font.bold = True

    y = 1.3
    for idx, lead in enumerate(briefing.top_leads[:8], start=1):
        acv_str = f"${lead.estimated_acv:,}" if lead.estimated_acv else "TBD"
        use_case_str = (lead.use_case or "?").replace("_", " ").title()
        company = lead.company_name or "Unknown"

        row = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11.8), Inches(0.55))
        tf = row.text_frame
        p = tf.paragraphs[0]
        p.text = f"{idx}. [{lead.grade}] {lead.score}  |  {company}  |  {use_case_str}  |  ACV: {acv_str}"
        p.font.size = Pt(13)
        if lead.grade in ("A", "B"):
            p.font.bold = True

        # Grade color bar
        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.2), Inches(y + 0.05),
            Inches(0.35), Inches(0.4),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = GRADE_COLORS.get(lead.grade, RGBColor(120, 120, 120))
        badge.line.fill.background()

        y += 0.7


def _add_lead_slide(prs: Presentation, lead: LeadCard) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    company = lead.company_name or "Unknown Company"

    # Company title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.5), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = company
    p.font.size = Pt(24)
    p.font.bold = True

    # Score badge
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.2), Inches(0.3), Inches(2.5), Inches(0.7),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GRADE_COLORS.get(lead.grade, RGBColor(120, 120, 120))
    badge.line.color.rgb = RGBColor(255, 255, 255)
    badge_tf = badge.text_frame
    acv_short = f"${lead.estimated_acv // 1000}K" if lead.estimated_acv and lead.estimated_acv >= 1000 else (f"${lead.estimated_acv:,}" if lead.estimated_acv else "TBD")
    badge_tf.text = f"{lead.grade} • {lead.score} • {acv_short}"
    badge_p = badge_tf.paragraphs[0]
    badge_p.alignment = PP_ALIGN.CENTER
    badge_p.font.bold = True
    badge_p.font.size = Pt(16)
    badge_p.font.color.rgb = RGBColor(255, 255, 255)

    # Contact + use case line with LinkedIn link
    contact = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12), Inches(0.4))
    p = contact.text_frame.paragraphs[0]
    title_part = f", {lead.sender_title}" if lead.sender_title else ""
    sender_part = lead.sender_name or "Unknown"
    use_case_str = (lead.use_case or "unknown").replace("_", " ").title()

    # Build LinkedIn search URL from name + company
    li_query = f"{sender_part} {lead.company_name or ''}".strip().replace(" ", "+")
    li_url = f"https://www.linkedin.com/search/results/all/?keywords={li_query}"

    # Contact info run
    run1 = p.add_run()
    run1.text = f"{sender_part}{title_part} — {lead.sender_email}  |  "
    run1.font.size = Pt(11)

    # LinkedIn hyperlink run
    run_li = p.add_run()
    run_li.text = "LinkedIn Profile"
    run_li.font.size = Pt(11)
    run_li.font.color.rgb = RGBColor(10, 102, 194)  # LinkedIn blue
    run_li.font.underline = True
    run_li.hyperlink.address = li_url

    # Use case run
    run3 = p.add_run()
    run3.text = f"  |  Use case: {use_case_str}"
    run3.font.size = Pt(11)

    # Left column: why hot + signals + MEDDPICC
    left = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.8), Inches(4.5))
    tf = left.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Why hot"
    p.font.bold = True
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = lead.why_hot
    p.font.size = Pt(12)

    p = tf.add_paragraph()
    p.text = "Key signals"
    p.font.bold = True
    p.font.size = Pt(14)

    for signal in lead.key_signals[:5]:
        p = tf.add_paragraph()
        p.text = f"• {signal}"
        p.font.size = Pt(11)

    if lead.meddpicc_summary:
        p = tf.add_paragraph()
        p.text = "MEDDPICC"
        p.font.bold = True
        p.font.size = Pt(14)

        p = tf.add_paragraph()
        p.text = lead.meddpicc_summary
        p.font.size = Pt(11)

    # Right column: context + competitive + approach
    right = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.5))
    tf = right.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Company context"
    p.font.bold = True
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = lead.company_context
    p.font.size = Pt(12)

    if lead.competitive_landscape:
        p = tf.add_paragraph()
        p.text = "Competitive landscape"
        p.font.bold = True
        p.font.size = Pt(14)

        p = tf.add_paragraph()
        p.text = lead.competitive_landscape
        p.font.size = Pt(12)

    p = tf.add_paragraph()
    p.text = "Suggested approach"
    p.font.bold = True
    p.font.size = Pt(14)

    p = tf.add_paragraph()
    p.text = lead.suggested_approach
    p.font.size = Pt(12)

    # Footer: opener + discovery question
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.15), Inches(12.0), Inches(1.0))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Suggested opener"
    p.font.bold = True
    p.font.size = Pt(13)

    p = tf.add_paragraph()
    p.text = lead.suggested_opener
    p.font.size = Pt(11)

    if lead.top_discovery_question:
        p = tf.add_paragraph()
        p.text = f"Top discovery question: {lead.top_discovery_question}"
        p.font.size = Pt(11)
        p.font.italic = True
