"""Call-prep endpoints — generate structured discovery call briefs via Claude."""

from __future__ import annotations

import json
import os
import sys
import tempfile
from datetime import date, timedelta
from pathlib import Path

from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel, field_validator

from server.db import get_lead, get_lead_emails, get_linkedin_one, get_research_one

# Allow agent imports
PROJECT_ROOT = str(Path(__file__).resolve().parent.parent.parent)
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

router = APIRouter(prefix="/api/leads", tags=["call-prep"])


def _is_gitlab_mode() -> bool:
    return os.environ.get("GITLAB_MODE", "false").lower() == "true"


# ── Coercion helpers ─────────────────────────────────────────────


def _coerce_str_list(v: object) -> list[str]:
    """Accept a list or a newline-delimited string."""
    if isinstance(v, list):
        return [str(x) for x in v]
    if isinstance(v, str):
        return [line for line in v.split("\n") if line.strip()]
    return []


def _coerce_objection_list(v: object) -> list[dict]:
    """Accept a list of dicts, a list of strings, or a single string."""
    if isinstance(v, list):
        out: list[dict] = []
        for item in v:
            if isinstance(item, dict):
                out.append({
                    "objection": str(item.get("objection", "")),
                    "response": str(item.get("response", "")),
                })
            else:
                out.append({"objection": str(item), "response": ""})
        return out
    if isinstance(v, str):
        return [{"objection": line, "response": ""} for line in v.split("\n") if line.strip()]
    return []


# ── Response schema ──────────────────────────────────────────────


class CallPrepResponse(BaseModel):
    contact_summary: str
    company_snapshot: str
    why_theyre_here: str
    current_stack: str
    suggested_agenda: list[str]
    discovery_questions: list[str]
    objection_prep: list[dict]  # [{objection, response}]
    competitive_positioning: str
    land_strategy: str
    expand_strategy: str
    proposed_next_step: str
    do_not_say: list[str]


class CallPrepExportRequest(BaseModel):
    contact_summary: str = ""
    company_snapshot: str = ""
    why_theyre_here: str = ""
    current_stack: str = ""
    suggested_agenda: list[str] = []
    discovery_questions: list[str] = []
    objection_prep: list[dict] = []
    competitive_positioning: str = ""
    land_strategy: str = ""
    expand_strategy: str = ""
    proposed_next_step: str = ""
    do_not_say: list[str] = []

    @field_validator("suggested_agenda", "discovery_questions", "do_not_say", mode="before")
    @classmethod
    def _coerce_str_lists(cls, v: object) -> list[str]:
        return _coerce_str_list(v)

    @field_validator("objection_prep", mode="before")
    @classmethod
    def _coerce_objections(cls, v: object) -> list[dict]:
        return _coerce_objection_list(v)


# Keep old name as alias so existing imports don't break
CallPrepPptxRequest = CallPrepExportRequest


# ── Helpers ──────────────────────────────────────────────────────


def _load_api_key() -> str:
    key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not key:
        from dotenv import load_dotenv
        load_dotenv(Path(PROJECT_ROOT) / ".env")
        key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not set")
    return key


def _gather_context(sender_email: str) -> dict:
    """Pull all available data for a lead from the database."""
    lead = get_lead(sender_email)
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    emails = get_lead_emails(sender_email)
    linkedin = get_linkedin_one(sender_email)

    # Fetch research by domain or company name
    research = None
    if lead.get("company_domain"):
        research = get_research_one(lead["company_domain"])
    if not research and lead.get("company_name"):
        company_key = lead["company_name"].lower().replace(" ", "_")
        research = get_research_one(company_key)

    return {
        "lead": lead,
        "emails": emails,
        "linkedin": linkedin,
        "research": research,
    }


def _build_prompt(ctx: dict) -> str:
    """Build the Claude prompt from all available context."""
    lead = ctx["lead"]
    linkedin = ctx["linkedin"]
    research = ctx["research"]
    emails = ctx["emails"]

    sections = []

    # Lead data
    sections.append(f"""## Lead Profile
- Name: {lead.get('sender_name') or 'Unknown'}
- Email: {lead['sender_email']}
- Title: {lead.get('sender_title') or 'Unknown'}
- Company: {lead.get('company_name') or 'Unknown'}
- Domain: {lead.get('company_domain') or 'Unknown'}
- Grade: {lead.get('grade')} ({lead.get('score')}/100)
- Intent: {lead.get('intent')}
- Use Case: {lead.get('use_case')}
- Product Fit: {lead.get('product_line_fit')}
- Buying Stage: {lead.get('buying_stage')}
- {'Current Toolchain' if _is_gitlab_mode() else 'Current AI Provider'}: {lead.get('current_ai_provider') or 'Unknown'}
- Estimated ACV: ${lead.get('estimated_acv') or 0:,}
- Deal Size Tier: {lead.get('deal_size_tier')}
- TAM Estimate: {lead.get('tam_estimate') or 'Unknown'}
- Expansion Potential: {lead.get('expansion_potential') or 'Unknown'}""")

    # Scoring breakdown
    sections.append(f"""## Scoring Breakdown
- Intent Score: {lead.get('intent_score')}/30
- Domain Quality: {lead.get('domain_quality_score')}/25
- Urgency: {lead.get('urgency_score')}/20
- Content Depth: {lead.get('content_depth_score')}/15
- Authority: {lead.get('authority_score')}/10
- Reasoning: {lead.get('last_reasoning') or 'None'}""")

    # MEDDPICC
    meddpicc = lead.get("meddpicc_json", {})
    if isinstance(meddpicc, dict) and meddpicc:
        meddpicc_lines = []
        for key, val in meddpicc.items():
            if isinstance(val, dict):
                status = val.get("status", "unknown")
                evidence = val.get("evidence", "")
                gap = val.get("gap", "")
                question = val.get("question", "")
                meddpicc_lines.append(
                    f"- {key}: status={status}, evidence={evidence}, gap={gap}, question={question}"
                )
            else:
                meddpicc_lines.append(f"- {key}: {val}")
        sections.append("## MEDDPICC Assessment\n" + "\n".join(meddpicc_lines))

    # Discovery questions
    disc_qs = lead.get("discovery_questions_json", [])
    if disc_qs:
        sections.append("## Existing Discovery Questions\n" + "\n".join(f"- {q}" for q in disc_qs))

    # Competitive signals
    comp_signals = lead.get("competitive_signals_json", [])
    if comp_signals:
        sections.append("## Competitive Signals\n" + "\n".join(f"- {s}" for s in comp_signals))

    # Partnership
    synergies = lead.get("partnership_synergies_json", [])
    detractors = lead.get("partnership_detractors_json", [])
    if synergies or detractors:
        parts = []
        if synergies:
            parts.append("Synergies:\n" + "\n".join(f"+ {s}" for s in synergies))
        if detractors:
            parts.append("Detractors:\n" + "\n".join(f"- {d}" for d in detractors))
        sections.append("## Partnership Analysis\n" + "\n".join(parts))

    # Next best action
    if lead.get("next_best_action"):
        sections.append(f"## Next Best Action\n{lead['next_best_action']}")

    # LinkedIn profile
    if linkedin:
        li_parts = [
            f"- Full Name: {linkedin.get('full_name') or 'Unknown'}",
            f"- Headline: {linkedin.get('headline') or 'Unknown'}",
            f"- Location: {linkedin.get('location') or 'Unknown'}",
            f"- Current Title: {linkedin.get('current_title') or 'Unknown'}",
            f"- Current Company: {linkedin.get('current_company') or 'Unknown'}",
        ]
        if linkedin.get("summary"):
            li_parts.append(f"- Summary: {linkedin['summary']}")

        exp = linkedin.get("experience_json", [])
        if exp:
            li_parts.append("Experience:")
            for e in exp[:5]:
                dur = f" ({e.get('duration', '')})" if e.get("duration") else ""
                li_parts.append(f"  - {e.get('title', '?')} at {e.get('company', '?')}{dur}")

        sections.append("## LinkedIn Profile\n" + "\n".join(li_parts))

    # Company research
    if research:
        res_parts = [
            f"- Company: {research.get('company_name')}",
            f"- Description: {research.get('description') or 'Unknown'}",
            f"- Employees: {research.get('employee_range') or 'Unknown'}",
            f"- Funding: {research.get('funding_stage') or 'Unknown'}",
        ]
        news = research.get("recent_news_json", [])
        if news:
            res_parts.append("Recent News:")
            for n in news[:5]:
                res_parts.append(f"  - {n.get('title', '')}: {n.get('summary', '')}")
        sections.append("## Company Research\n" + "\n".join(res_parts))

    # Emails
    if emails:
        email_parts = []
        for e in emails[:5]:
            email_parts.append(
                f"Subject: {e.get('subject', '')}\n"
                f"Date: {e.get('received_at', '')}\n"
                f"Body:\n{e.get('body', '')}\n---"
            )
        sections.append("## Email History\n" + "\n".join(email_parts))

    return "\n\n".join(sections)


# ── Endpoints ────────────────────────────────────────────────────


CALL_PREP_TOOL = {
    "name": "generate_call_prep",
    "description": "Generate a structured call preparation brief for a discovery call.",
    "input_schema": {
        "type": "object",
        "required": [
            "contact_summary", "company_snapshot", "why_theyre_here",
            "current_stack", "suggested_agenda", "discovery_questions",
            "objection_prep", "competitive_positioning", "land_strategy",
            "expand_strategy", "proposed_next_step", "do_not_say",
        ],
        "properties": {
            "contact_summary": {
                "type": "string",
                "description": "2-3 sentences: who you're talking to — name, title, background from LinkedIn, what they likely care about based on their inbound email and role.",
            },
            "company_snapshot": {
                "type": "string",
                "description": "One paragraph: company overview including estimated revenue, employees, funding stage, recent news, and relevant tech stack details.",
            },
            "why_theyre_here": {
                "type": "string",
                "description": "The specific pain points and motivations extracted from their inbound email. What triggered them to reach out.",
            },
            "current_stack": {
                "type": "string",
                "description": "What AI provider they use now and what's working/not working based on competitive signals and email content.",
            },
            "suggested_agenda": {
                "type": "array",
                "items": {"type": "string"},
                "description": "5-bullet meeting framework for a 25-minute discovery call.",
            },
            "discovery_questions": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Top 5 questions tailored to this specific lead. Refine and improve the existing discovery questions based on all available context.",
            },
            "objection_prep": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "objection": {"type": "string"},
                        "response": {"type": "string"},
                    },
                    "required": ["objection", "response"],
                },
                "description": "2-3 likely objections based on their current provider and situation, with suggested responses.",
            },
            "competitive_positioning": {
                "type": "string",
                "description": "The product's specific advantages vs their current provider for their specific use cases.",
            },
            "land_strategy": {
                "type": "string",
                "description": "What to propose as the initial deal — specific product, number of seats, and pricing approach.",
            },
            "expand_strategy": {
                "type": "string",
                "description": "The full account TAM and how to grow the deal over 12 months across departments and products.",
            },
            "proposed_next_step": {
                "type": "string",
                "description": "What to propose at the end of the call as the concrete next step.",
            },
            "do_not_say": {
                "type": "array",
                "items": {"type": "string"},
                "description": "Things to avoid mentioning — e.g. don't trash their existing integration if they just shipped it, don't mention features that aren't GA yet.",
            },
        },
    },
}


@router.post("/{sender_email:path}/call-prep")
async def generate_call_prep(sender_email: str):
    """Generate a structured call-prep brief using Claude."""
    api_key = _load_api_key()
    ctx = _gather_context(sender_email)
    prompt = _build_prompt(ctx)

    import anthropic

    client = anthropic.Anthropic(api_key=api_key)

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4096,
        tools=[CALL_PREP_TOOL],
        tool_choice={"type": "tool", "name": "generate_call_prep"},
        messages=[
            {
                "role": "user",
                "content": (
                    "You are preparing a sales rep for a 25-minute discovery call. "
                    "Generate a comprehensive call preparation brief based on all the data below. "
                    "Be specific and actionable — reference actual details from the emails, LinkedIn profile, "
                    "and company research. Do not be generic.\n\n"
                    f"{prompt}"
                ),
            }
        ],
    )

    # Extract tool use result
    for block in response.content:
        if block.type == "tool_use" and block.name == "generate_call_prep":
            return block.input

    raise HTTPException(status_code=500, detail="Claude did not return a call-prep result")


# ── A) Pre-Call Brief PDF ────────────────────────────────────────


@router.post("/{sender_email:path}/call-prep-pdf")
async def generate_call_prep_pdf(sender_email: str, data: CallPrepExportRequest):
    """Generate a dense single-page PDF reference brief."""
    lead = get_lead(sender_email)
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    output_path = _build_call_prep_pdf(lead, data)
    return FileResponse(
        path=str(output_path),
        media_type="application/pdf",
        filename=f"pre-call-brief-{lead.get('company_name', 'prospect').replace(' ', '-').lower()}.pdf",
    )


def _sanitize(text: str) -> str:
    """Replace Unicode characters that Helvetica can't render."""
    return (
        text
        .replace("\u2022", "-")   # •
        .replace("\u2013", "-")   # –
        .replace("\u2014", "--")  # —
        .replace("\u2018", "'")   # '
        .replace("\u2019", "'")   # '
        .replace("\u201c", '"')   # "
        .replace("\u201d", '"')   # "
        .replace("\u2026", "...")  # …
        .replace("\u00a0", " ")   # non-breaking space
        .replace("\u200b", "")    # zero-width space
    )


def _build_call_prep_pdf(lead: dict, data: CallPrepExportRequest) -> Path:
    from fpdf import FPDF

    company = lead.get("company_name") or "Prospect"
    contact = lead.get("sender_name") or "Contact"
    title_str = lead.get("sender_title") or ""
    grade = lead.get("grade", "?")
    score = lead.get("score", 0)
    acv = lead.get("estimated_acv")
    acv_str = f"${acv:,}" if acv else "TBD"
    today = date.today().strftime("%B %d, %Y")

    pdf = FPDF(orientation="L", unit="mm", format="letter")
    pdf.set_auto_page_break(auto=False)
    W = 279.4  # letter landscape width
    H = 215.9  # letter landscape height
    M = 8      # margin

    # ── Layout constants ──
    BODY_FONT_SIZE = 7.5
    HEADER_FONT_SIZE = 8
    COL_W = (W - M * 3) / 2
    LEFT_X = M
    RIGHT_X = M + COL_W + M
    START_Y = 22
    LINE_H = 4.2
    SECTION_GAP = 3.0
    MAX_Y = H - 8  # printable bottom

    def _draw_header(page_label: str) -> None:
        """Draw dark header bar with accent line on the current page."""
        pdf.set_fill_color(26, 26, 26)
        pdf.rect(0, 0, W, 18, "F")
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_text_color(255, 255, 255)
        pdf.set_xy(M, 3)
        pdf.cell(0, 5, f"PRE-CALL BRIEF: {company.upper()}", new_x="LMARGIN")
        pdf.set_font("Helvetica", "", 7)
        pdf.set_xy(W - M - 40, 3)
        pdf.cell(40, 5, page_label, align="R", new_x="LMARGIN")
        pdf.set_font("Helvetica", "", 8.5)
        pdf.set_xy(M, 10)
        detail = f"{contact}"
        if title_str:
            detail += f", {title_str}"
        detail += f"   |   Grade {grade}  ({score}/100)   |   ACV {acv_str}   |   {today}"
        pdf.cell(0, 4, detail, new_x="LMARGIN")
        if _is_gitlab_mode():
            pdf.set_draw_color(252, 109, 38)
        else:
            pdf.set_draw_color(212, 165, 116)
        pdf.set_line_width(0.8)
        pdf.line(0, 18, W, 18)

    def _section_header(x: float, y: float, text: str) -> float:
        pdf.set_font("Helvetica", "B", HEADER_FONT_SIZE)
        if _is_gitlab_mode():
            pdf.set_text_color(252, 109, 38)
        else:
            pdf.set_text_color(212, 165, 116)
        pdf.set_xy(x, y)
        pdf.cell(COL_W, LINE_H, text.upper(), new_x="LMARGIN")
        pdf.set_draw_color(200, 200, 200)
        pdf.set_line_width(0.15)
        pdf.line(x, y + LINE_H + 0.3, x + COL_W, y + LINE_H + 0.3)
        return y + LINE_H + 1.5

    def _body_text(x: float, y: float, text: str) -> float:
        pdf.set_font("Helvetica", "", BODY_FONT_SIZE)
        pdf.set_text_color(50, 50, 50)
        pdf.set_xy(x, y)
        pdf.multi_cell(COL_W, LINE_H, _sanitize(text.replace("\n", " ").strip()), new_x="LMARGIN")
        return pdf.get_y() + SECTION_GAP

    def _bullet_list(x: float, y: float, items: list[str], max_items: int = 6) -> float:
        pdf.set_font("Helvetica", "", BODY_FONT_SIZE)
        pdf.set_text_color(50, 50, 50)
        for i, item in enumerate(items[:max_items]):
            pdf.set_xy(x, y)
            pdf.multi_cell(COL_W, LINE_H, _sanitize(f"{i + 1}. {item.strip()}"), new_x="LMARGIN")
            y = pdf.get_y() + 0.5
        return y + SECTION_GAP

    # ════════════════════════════════════════════════════════════════
    # PAGE 1 — INTELLIGENCE (context before the call)
    # ════════════════════════════════════════════════════════════════
    pdf.add_page()
    _draw_header("PAGE 1: INTELLIGENCE")

    # ── Left column ──
    left_y = START_Y

    left_y = _section_header(LEFT_X, left_y, "Contact Summary")
    left_y = _body_text(LEFT_X, left_y, data.contact_summary)

    left_y = _section_header(LEFT_X, left_y, "Company Snapshot")
    left_y = _body_text(LEFT_X, left_y, data.company_snapshot)

    left_y = _section_header(LEFT_X, left_y, "Why They're Here")
    left_y = _body_text(LEFT_X, left_y, data.why_theyre_here)

    left_y = _section_header(LEFT_X, left_y, "Current Stack")
    left_y = _body_text(LEFT_X, left_y, data.current_stack)

    # ── Right column ──
    right_y = START_Y

    right_y = _section_header(RIGHT_X, right_y, "Competitive Positioning")
    right_y = _body_text(RIGHT_X, right_y, data.competitive_positioning)

    right_y = _section_header(RIGHT_X, right_y, "Suggested Agenda")
    right_y = _bullet_list(RIGHT_X, right_y, data.suggested_agenda, max_items=5)

    # Confidential line at bottom of page 1
    pdf.set_font("Helvetica", "I", 5)
    pdf.set_text_color(150, 150, 150)
    pdf.set_xy(M, H - 6)
    pdf.cell(W - M * 2, 3, "INTERNAL USE ONLY  |  Blackburn Command Center  |  Generated by Claude", new_x="LMARGIN")

    # ════════════════════════════════════════════════════════════════
    # PAGE 2 — ACTION PLAN (what to do on the call)
    # ════════════════════════════════════════════════════════════════
    pdf.add_page()
    _draw_header("PAGE 2: ACTION PLAN")

    # ── Left column ──
    left_y = START_Y

    left_y = _section_header(LEFT_X, left_y, "Discovery Questions")
    left_y = _bullet_list(LEFT_X, left_y, data.discovery_questions, max_items=6)

    left_y = _section_header(LEFT_X, left_y, "Proposed Next Step")
    pdf.set_font("Helvetica", "B", BODY_FONT_SIZE)
    pdf.set_text_color(26, 26, 26)
    pdf.set_xy(LEFT_X, left_y)
    pdf.multi_cell(COL_W, LINE_H, _sanitize(data.proposed_next_step.strip()), new_x="LMARGIN")
    left_y = pdf.get_y() + SECTION_GAP

    # ── Right column ──
    right_y = START_Y

    right_y = _section_header(RIGHT_X, right_y, "Objection Prep")
    for obj in data.objection_prep[:3]:
        objection = obj.get("objection", "")
        response = obj.get("response", "")
        pdf.set_text_color(180, 60, 60)
        pdf.set_font("Helvetica", "B", 7)
        pdf.set_xy(RIGHT_X, right_y)
        pdf.multi_cell(COL_W, LINE_H, _sanitize(f"  {objection}"), new_x="LMARGIN")
        right_y = pdf.get_y()
        pdf.set_text_color(50, 50, 50)
        pdf.set_font("Helvetica", "", 7)
        pdf.set_xy(RIGHT_X + 3, right_y)
        pdf.multi_cell(COL_W - 3, LINE_H, _sanitize(response), new_x="LMARGIN")
        right_y = pdf.get_y() + 1.5
    right_y += SECTION_GAP

    right_y = _section_header(RIGHT_X, right_y, "Deal Strategy")
    pdf.set_font("Helvetica", "", BODY_FONT_SIZE)
    pdf.set_text_color(50, 50, 50)
    pdf.set_xy(RIGHT_X, right_y)
    pdf.multi_cell(COL_W, LINE_H, _sanitize(f"Land: {data.land_strategy.strip()}"), new_x="LMARGIN")
    right_y = pdf.get_y() + 1
    pdf.set_xy(RIGHT_X, right_y)
    pdf.multi_cell(COL_W, LINE_H, _sanitize(f"Expand: {data.expand_strategy.strip()}"), new_x="LMARGIN")
    right_y = pdf.get_y() + SECTION_GAP

    # ── Footer: Do Not Say ──
    footer_y = max(left_y, right_y, H - 30) + 2
    if footer_y > MAX_Y - 15:
        footer_y = MAX_Y - 15

    pdf.set_draw_color(200, 200, 200)
    pdf.set_line_width(0.15)
    pdf.line(M, footer_y, W - M, footer_y)
    footer_y += 1.5

    dns_items = data.do_not_say[:5]
    if dns_items:
        pdf.set_font("Helvetica", "BI", 6.5)
        pdf.set_text_color(180, 60, 60)
        pdf.set_xy(M, footer_y)
        pdf.cell(W - M * 2, LINE_H, "DO NOT SAY:", new_x="LMARGIN")
        footer_y += LINE_H
        pdf.set_font("Helvetica", "I", 6)
        for item in dns_items:
            pdf.set_xy(M + 3, footer_y)
            pdf.multi_cell(W - M * 2 - 3, LINE_H, _sanitize(f"- {item}"), new_x="LMARGIN")
            footer_y = pdf.get_y()

    # Confidential footer
    pdf.set_font("Helvetica", "I", 5)
    pdf.set_text_color(150, 150, 150)
    pdf.set_xy(M, H - 6)
    pdf.cell(W - M * 2, 3, "INTERNAL USE ONLY  |  Blackburn Command Center  |  Generated by Claude", new_x="LMARGIN")

    tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    pdf.output(tmp.name)
    return Path(tmp.name)


# ── B) Customer-Facing PPTX Deck ────────────────────────────────


@router.post("/{sender_email:path}/call-prep-pptx")
async def generate_call_prep_pptx(sender_email: str, data: CallPrepExportRequest):
    """Generate a branded customer-facing PPTX deck with speaker notes."""
    lead = get_lead(sender_email)
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")

    output_path = _build_customer_deck(lead, data)
    return FileResponse(
        path=str(output_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{'gitlab' if _is_gitlab_mode() else 'anthropic'}-{lead.get('company_name', 'prospect').replace(' ', '-').lower()}.pptx",
    )


def _build_customer_deck(lead: dict, data: CallPrepExportRequest) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # ── Brand colours ──
    if _is_gitlab_mode():
        BG_CREAM = RGBColor(23, 19, 33)       # #171321
        TEXT_DARK = RGBColor(255, 255, 255)
        TEXT_BODY = RGBColor(200, 200, 200)
        ACCENT = RGBColor(252, 109, 38)        # #FC6D26
        MUTED = RGBColor(162, 161, 166)
        WHITE = RGBColor(255, 255, 255)
    else:
        BG_CREAM = RGBColor(245, 240, 232)
        TEXT_DARK = RGBColor(26, 26, 26)
        TEXT_BODY = RGBColor(68, 64, 60)
        ACCENT = RGBColor(212, 165, 116)
        MUTED = RGBColor(140, 135, 125)
        WHITE = RGBColor(255, 255, 255)
    FONT = "Arial"

    company = lead.get("company_name") or "Your Company"
    contact = lead.get("sender_name") or "Contact"
    today_str = date.today().strftime("%B %d, %Y")
    acv = lead.get("estimated_acv")
    acv_str = f"${acv:,}" if acv else "TBD"
    emp_count = lead.get("employee_range") or "793"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Primitives ──

    def _bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_CREAM

    def _t(slide, x, y, w, h, text, size=16, bold=False,
           color=TEXT_DARK, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT
        p.alignment = align
        return tf

    def _n(slide, text: str):
        slide.notes_slide.notes_text_frame.text = text

    def _bar(slide, x, y, w, h):
        s = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = ACCENT
        s.line.fill.background()

    def _card(slide, x, y, w, h, fill=WHITE):
        s = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.color.rgb = RGBColor(220, 215, 205)
        s.line.width = Pt(0.5)

    _brand_name = "gitlab" if _is_gitlab_mode() else "anthropic"

    def _header(slide, title: str):
        _bar(slide, 0, 0, 13.333, 0.08)
        _t(slide, 1.0, 0.6, 6, 0.5, _brand_name, size=14, color=ACCENT, bold=True)
        _t(slide, 1.0, 1.3, 10, 0.8, title, size=36, bold=True)
        _bar(slide, 1.0, 2.2, 1.0, 0.05)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s1)
    _bar(s1, 0, 0, 13.333, 0.08)
    _t(s1, 1.0, 0.6, 6, 0.5, _brand_name, size=14, color=ACCENT, bold=True)
    _t(s1, 1.0, 2.0, 10, 1.5,
       f"Partnering with {company}\non AI", size=40, bold=True)
    _bar(s1, 1.0, 4.0, 1.5, 0.06)
    _t(s1, 1.0, 4.4, 8, 0.5,
       f"Prepared for {contact}  |  {today_str}", size=16, color=MUTED)

    _n(s1, (
        f"INTRO TALK TRACK\n\n"
        f"Thank {contact} for taking the time. Reference their inbound interest.\n\n"
        f"CONTACT SUMMARY:\n{data.contact_summary}\n\n"
        f"COMPANY BACKGROUND:\n{data.company_snapshot}\n\n"
        f"Key: Establish credibility, show you've done homework, set collaborative tone. "
        f"This is a conversation, not a pitch."
    ))

    # ════════════════════════════════════════════════════════════════
    # SLIDE 2 — Agenda (hardcoded short labels)
    # ════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s2)
    _header(s2, "Today's Agenda")

    agenda_labels = [
        "Your AI Landscape Today",
        "Key Challenges & Requirements",
        "How GitLab Can Help" if _is_gitlab_mode() else "How Claude Can Help",
        "Pilot Approach & Next Steps",
    ]
    y = 2.6
    for label in agenda_labels:
        _card(s2, 1.0, y, 11.0, 0.9)
        _t(s2, 1.5, y + 0.22, 10.0, 0.5, label, size=18, color=TEXT_BODY)
        y += 1.05

    agenda_items = data.suggested_agenda[:5]
    time_allocs = ["3 min", "5 min", "7 min", "5 min", "5 min"]
    if len(agenda_items) == 4:
        time_allocs = ["3 min", "7 min", "8 min", "7 min"]
    notes_agenda = "AGENDA -- FULL DETAIL & TIME ALLOCATIONS\n\n"
    for i, item in enumerate(agenda_items):
        t = time_allocs[i] if i < len(time_allocs) else "5 min"
        notes_agenda += f"{i + 1}. ({t}) {item}\n"
    notes_agenda += (
        f"\nHow to frame: 'I want to make sure we use our time well. "
        f"I'd love to start by understanding your current landscape, "
        f"then share how we might be able to help, and end with clear next steps. "
        f"Does that work?'"
    )
    _n(s2, notes_agenda)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 3 — Understanding Your AI Landscape (hardcoded cards)
    # ════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s3)
    _header(s3, "Understanding Your AI Landscape")

    challenges = [
        ("Accuracy at Stake",
         "Contract errors putting enterprise clients at risk"),
        ("Governance Gap",
         f"{emp_count} employees, no centralized AI oversight"),
        ("Legal-Grade Reliability",
         "Fortune 500 contracts demand zero tolerance"),
    ]
    card_xs = [1.0, 4.8, 8.6]
    for i, (headline, subtitle) in enumerate(challenges):
        _card(s3, card_xs[i], 2.9, 3.5, 3.2)
        _bar(s3, card_xs[i], 2.9, 3.5, 0.06)
        _t(s3, card_xs[i] + 0.3, 3.3, 2.9, 0.5,
           f"0{i + 1}", size=32, bold=True, color=ACCENT)
        _t(s3, card_xs[i] + 0.3, 4.0, 2.9, 0.6,
           headline, size=22, bold=True)
        _t(s3, card_xs[i] + 0.3, 4.7, 2.9, 1.2,
           subtitle, size=14, color=TEXT_BODY)

    notes_landscape = (
        f"DISCOVERY SECTION -- FULL DETAIL\n\n"
        f"WHY THEY'RE HERE:\n{data.why_theyre_here}\n\n"
        f"CURRENT STACK:\n{data.current_stack}\n\n"
        f"DISCOVERY QUESTIONS TO ASK:\n"
    )
    for i, q in enumerate(data.discovery_questions[:6], 1):
        notes_landscape += f"{i}. {q}\n"
    notes_landscape += (
        f"\nListen more than you talk in this section. Use these questions "
        f"to uncover the real pain. Take notes on what they say -- you'll "
        f"reference it in the next section.\n\n"
        f"If they mention specific numbers (error counts, ticket volumes, "
        f"cost figures), write them down verbatim. You'll use these in the "
        f"proposal section."
    )
    _n(s3, notes_landscape)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 4 — Why Claude for {company} (hardcoded value props)
    # ════════════════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s4)
    _header(s4, f"Why {'GitLab' if _is_gitlab_mode() else 'Claude'} for {company}")

    value_props = [
        ("Superior Legal Reasoning",
         "Constitutional AI built for nuanced interpretation"),
        ("Enterprise-Ready from Day One",
         "SSO, security controls, and compliance built in"),
        ("Seamless Migration Path",
         "API-compatible with dedicated migration support"),
    ]
    y = 2.9
    for headline, subtitle in value_props:
        _card(s4, 1.0, y, 11.0, 1.35)
        _bar(s4, 1.0, y, 0.08, 1.35)
        _t(s4, 1.5, y + 0.2, 10.0, 0.6,
           headline, size=22, bold=True)
        _t(s4, 1.5, y + 0.75, 10.0, 0.5,
           subtitle, size=14, color=TEXT_BODY)
        y += 1.55

    notes_why = (
        f"COMPETITIVE POSITIONING -- FULL DETAIL\n\n"
        f"{data.competitive_positioning}\n\n"
        f"OBJECTION RESPONSES (word-for-word rebuttals):\n"
    )
    for obj in data.objection_prep[:4]:
        notes_why += f"\nOBJECTION: \"{obj.get('objection', '')}\"\n"
        notes_why += f"RESPONSE: {obj.get('response', '')}\n"
    notes_why += (
        f"\nKey: Connect advantages to THEIR specific pain points. "
        f"Reference what they said in the discovery section. "
        f"'Based on what you shared about X, here\\'s why {'GitLab' if _is_gitlab_mode() else 'Claude'} is a good fit...'\n\n"
        f"Do NOT lead with features. Lead with their pain, then bridge to how {'GitLab' if _is_gitlab_mode() else 'Claude'} solves it."
    )
    _n(s4, notes_why)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 5 — Proposed Pilot Program (no dollar amounts on slide)
    # ════════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s5)
    _header(s5, "Proposed Pilot Program")

    phases = [
        ("PHASE 1", "30 Days", "Technical Proof-of-Concept"),
        ("PHASE 2", "90 Days", "Expand to Key Teams"),
        ("PHASE 3", "12 Months", "Full Organization Rollout"),
    ]

    # Connector line
    conn = s5.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1.8), Inches(3.7), Inches(9.6), Inches(0.06))
    conn.fill.solid()
    conn.fill.fore_color.rgb = ACCENT
    conn.line.fill.background()

    phase_xs = [1.2, 5.0, 8.8]
    for i, (label, time, desc) in enumerate(phases):
        circle = s5.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(phase_xs[i] + 1.3), Inches(3.4), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if i == 0 else BG_CREAM
        circle.line.color.rgb = ACCENT
        circle.line.width = Pt(2)
        _t(s5, phase_xs[i] + 1.3, 3.45, 0.7, 0.6,
           str(i + 1), size=22, bold=True,
           color=WHITE if i == 0 else ACCENT,
           align=PP_ALIGN.CENTER)

        _card(s5, phase_xs[i], 4.4, 3.6, 2.0)
        _t(s5, phase_xs[i] + 0.25, 4.55, 3.1, 0.4,
           f"{label}  --  {time}", size=14, bold=True, color=ACCENT)
        _t(s5, phase_xs[i] + 0.25, 5.1, 3.1, 0.7,
           desc, size=16, color=TEXT_BODY)

    notes_pilot = (
        f"FULL DEAL STRATEGY\n\n"
        f"LAND STRATEGY:\n{data.land_strategy}\n\n"
        f"EXPAND STRATEGY (12-MONTH):\n{data.expand_strategy}\n\n"
        f"ACV Target: {acv_str}\n"
        f"TAM Estimate: {lead.get('tam_estimate') or 'TBD'}\n"
        f"Expansion Potential: {lead.get('expansion_potential') or 'TBD'}\n\n"
        f"Present this as a low-risk starting point. Emphasize time-to-value.\n"
        f"Phase 1 should feel easy to say yes to.\n\n"
        f"If they push back on timeline, offer to compress Phase 1 to 2 weeks "
        f"with a focused use case. The goal is to get code running against "
        f"their data as fast as possible."
    )
    _n(s5, notes_pilot)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 6 — Next Steps (hardcoded short labels)
    # ════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s6)
    _header(s6, "Next Steps")

    d1 = (date.today() + timedelta(days=2)).strftime("%b %d")
    d2 = (date.today() + timedelta(days=7)).strftime("%b %d")
    d3 = (date.today() + timedelta(days=14)).strftime("%b %d")

    actions = [
        ("Board presentation materials", "Parker / GitLab" if _is_gitlab_mode() else "Parker / Anthropic", d1),
        ("Technical deep-dive session", contact, d2),
        ("Pilot kickoff", "Joint", d3),
    ]

    # Table header
    _card(s6, 1.0, 2.8, 11.0, 0.65, fill=RGBColor(235, 230, 222))
    _t(s6, 1.4, 2.9, 6, 0.45, "Action Item", size=14, bold=True, color=MUTED)
    _t(s6, 8.0, 2.9, 2, 0.45, "Owner", size=14, bold=True, color=MUTED)
    _t(s6, 10.5, 2.9, 2, 0.45, "Target Date", size=14, bold=True, color=MUTED)

    y = 3.6
    for action, owner, dt in actions:
        _card(s6, 1.0, y, 11.0, 0.85)
        _t(s6, 1.4, y + 0.2, 6.2, 0.5, action, size=16, color=TEXT_BODY)
        _t(s6, 8.0, y + 0.2, 2.2, 0.5, owner, size=14, color=TEXT_BODY)
        _t(s6, 10.5, y + 0.2, 2, 0.5, dt, size=14, bold=True, color=ACCENT)
        y += 1.0

    _t(s6, 1.0, 6.5, 10, 0.5,
       "Parker Blackburn  |  parker@gitlab.com  |  GitLab" if _is_gitlab_mode() else "Parker Blackburn  |  parker@anthropic.com  |  Anthropic",
       size=14, color=MUTED)

    notes_next = (
        f"CLOSING SECTION -- FULL DETAIL\n\n"
        f"PROPOSED NEXT STEP:\n{data.proposed_next_step}\n\n"
        f"Push for a specific date and time. Don't leave with 'we'll follow up.'\n"
        f"Ideal outcome: calendar invite sent before the call ends.\n\n"
        f"DO NOT SAY:\n"
    )
    for item in data.do_not_say[:5]:
        notes_next += f"- {item}\n"
    notes_next += (
        f"\nClosing language: 'Based on everything we've discussed, I think "
        f"the best next step would be [proposed step]. Can we get that on the "
        f"calendar before we wrap up today?'"
    )
    _n(s6, notes_next)

    # ── Save ──
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return Path(tmp.name)
